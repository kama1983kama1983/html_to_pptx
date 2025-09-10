import os
import zipfile
from pathlib import Path
from flask import Flask, render_template_string, request, send_file, url_for, jsonify
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches
from playwright.sync_api import sync_playwright
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader
from PIL import Image

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


# ----------------------------
# Convert HTML to PNG
# ----------------------------
def html_to_png(html_path, output_png_path, viewport={"width": 1280, "height": 720}):
    print(f"[🔄] تحويل {html_path.name} → PNG")
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport=viewport)
        file_url = f"file://{html_path.resolve()}"
        page.goto(file_url, wait_until="networkidle")
        page.screenshot(path=str(output_png_path), full_page=True)
        browser.close()
    print(f"[✅] تم إنشاء {output_png_path.name}")


# ----------------------------
# Convert list of PNGs to PPTX
# ----------------------------
def images_to_pptx(image_files, pptx_path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for img_path in image_files:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str(img_path), Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)
    prs.save(pptx_path)
    print(f"[✅] PPTX محفوظ: {pptx_path}")


# ----------------------------
# Convert list of PNGs to PDF
# ----------------------------
def images_to_pdf(image_files, pdf_path):
    c = canvas.Canvas(str(pdf_path), pagesize=A4)
    width, height = A4
    for img_path in image_files:
        img = Image.open(img_path)
        img.thumbnail((width, height))
        c.drawImage(ImageReader(img), 0, 0, width=width, height=height)
        c.showPage()
    c.save()
    print(f"[✅] PDF محفوظ: {pdf_path}")


# ----------------------------
# Flask Routes
# ----------------------------
@app.route("/", methods=["GET", "POST"])
def upload():
    if request.method == "POST":
        file = request.files["file"]
        sort_mode = request.form.get("sort", "name")  # name أو date

        if file and file.filename.endswith(".zip"):
            filename = secure_filename(file.filename)
            zip_path = Path(UPLOAD_FOLDER) / filename
            file.save(zip_path)
            print(f"[⬆️] رفع الملف: {filename}")

            # Extract zip
            extract_dir = Path(UPLOAD_FOLDER) / filename.replace(".zip", "")
            os.makedirs(extract_dir, exist_ok=True)
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)

            # ترتيب الملفات
            if sort_mode == "date":
                html_files = sorted(extract_dir.glob("*.html"),
                                    key=lambda f: f.stat().st_mtime)
            else:  # default بالاسم
                html_files = sorted(extract_dir.glob("*.html"),
                                    key=lambda f: f.name)

            print(f"[🔍] وجدنا {len(html_files)} ملف HTML (ترتيب: {sort_mode})")

            # تحويل HTML → PNG
            out_dir = Path(OUTPUT_FOLDER) / filename.replace(".zip", "")
            os.makedirs(out_dir, exist_ok=True)
            image_files = []
            for html_file in html_files:
                out_png = out_dir / (html_file.stem + ".png")
                html_to_png(html_file, out_png)
                image_files.append(out_png)

            # PPTX
            pptx_path = out_dir / "presentation.pptx"
            images_to_pptx(image_files, pptx_path)

            # PDF
            pdf_path = out_dir / "presentation.pdf"
            images_to_pdf(image_files, pdf_path)

            return render_template_string("""
                <!doctype html>
                <html lang="ar">
                <head>
                    <meta charset="utf-8">
                    <title>مراجعة الشرائح</title>
                    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css" rel="stylesheet">
                </head>
                <body class="container py-4">
                    <h2>✅ تم إنشاء العرض التقديمي</h2>
                    <a href="{{ url_for('download', filename=pptx_file) }}" class="btn btn-success mb-2">⬇️ تحميل PPTX</a>
                    <a href="{{ url_for('download', filename=pdf_file) }}" class="btn btn-danger mb-2">⬇️ تحميل PDF</a>
                    
                    <h3 class="mt-4">📑 مراجعة:</h3>
                    <ul class="list-group">
                        {% for html in htmls %}
                        <li class="list-group-item">
                            صفحة {{ loop.index }}: {{ html }}
                            <iframe src="{{ url_for('serve_upload', filename=html) }}" width="100%" height="200"></iframe>
                        </li>
                        {% endfor %}
                    </ul>
                </body>
                </html>
            """, pptx_file=str(Path(filename.replace(".zip", "")) / "presentation.pptx"),
                 pdf_file=str(Path(filename.replace(".zip", "")) / "presentation.pdf"),
                 htmls=[str(Path(filename.replace(".zip", "")) / f.name) for f in html_files])

    return """
        <!doctype html>
        <html lang="ar">
        <head>
            <meta charset="utf-8">
            <title>رفع ملفات HTML</title>
            <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css" rel="stylesheet">
        </head>
        <body class="container py-5">
            <h1>📤 رفع ملفات HTML مضغوطة</h1>
            <form method="post" enctype="multipart/form-data" class="mb-3">
                <input type="file" name="file" accept=".zip" class="form-control mb-3">
                <label class="form-label">ترتيب الملفات:</label>
                <select name="sort" class="form-select mb-3">
                    <option value="name">حسب الاسم (page1, page2...)</option>
                    <option value="date">حسب تاريخ الإنشاء</option>
                </select>
                <button type="submit" class="btn btn-primary">رفع ومعالجة</button>
            </form>
        </body>
        </html>
    """


@app.route("/download/<path:filename>")
def download(filename):
    return send_file(Path(OUTPUT_FOLDER) / filename, as_attachment=True)


@app.route("/uploads/<path:filename>")
def serve_upload(filename):
    return send_file(Path(UPLOAD_FOLDER) / filename)


if __name__ == "__main__":
    print("[🚀] التطبيق يعمل على http://127.0.0.1:5000")
    app.run(debug=True, use_reloader=False)
